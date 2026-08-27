"""Stage 'deck': 10-slide executive PowerPoint built from the report tables and figures."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from supply_pipeline.config import Config

log = logging.getLogger(__name__)

INK = RGBColor(0x18, 0x22, 0x1D)
MUTED = RGBColor(0x5B, 0x68, 0x62)
ACCENT = RGBColor(0x0E, 0x6B, 0x58)
GOLD = RGBColor(0xB9, 0x91, 0x1E)
LINE = RGBColor(0xD6, 0xDD, 0xD6)
W, H = Inches(13.333), Inches(7.5)
STAT_WIDTH = Inches(2.6)
FLOW_BOX_HEIGHT = Inches(1.1)


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self.blank = self.prs.slide_layouts[6]
        self.n = 0

    def slide(self, title: str, kicker: str | None = None) -> Any:
        s = self.prs.slides.add_slide(self.blank)
        self.n += 1
        if kicker:
            self.text(s, kicker, Inches(0.6), Inches(0.35), Inches(12), Inches(0.3), size=11, color=ACCENT, bold=True)
        self.text(s, title, Inches(0.6), Inches(0.6), Inches(12), Inches(0.8), size=26, bold=True)
        ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.6), Inches(1.4), Inches(12.7), Inches(1.4))
        ln.line.color.rgb = LINE
        self.text(
            s, f"{self.n} / 10", Inches(12.2), Inches(7.0), Inches(0.9), Inches(0.3), size=9, color=MUTED, align=PP_ALIGN.RIGHT
        )
        return s

    @staticmethod
    def text(
        s: Any,
        txt: str,
        x: int,
        y: int,
        w: int,
        h: int,
        size: int = 14,
        bold: bool = False,
        color: RGBColor = INK,
        align: Any = PP_ALIGN.LEFT,
    ) -> Any:
        tb = s.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = txt
        p.alignment = align
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        return tb

    @staticmethod
    def bullets(s: Any, items: list[str], x: int, y: int, w: int, h: int, size: int = 13) -> Any:
        tb = s.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "•  " + it
            p.font.size = Pt(size)
            p.font.color.rgb = INK
            p.space_after = Pt(6)
        return tb

    @staticmethod
    def picture(s: Any, path: Path, x: int, y: int, w: int | None = None, h: int | None = None) -> Any:
        if not path.exists():
            return None
        return s.shapes.add_picture(str(path), x, y, width=w, height=h)

    def table(self, s: Any, df: pd.DataFrame, x: int, y: int, w: int, h: int, size: int = 10, fmt: str = "{:.2f}") -> Any:
        rows, cols = df.shape[0] + 1, df.shape[1]
        tbl = s.shapes.add_table(rows, cols, x, y, w, h).table
        for j, c in enumerate(df.columns):
            cell = tbl.cell(0, j)
            cell.text = str(c)
            cell.text_frame.paragraphs[0].font.size = Pt(size)
            cell.text_frame.paragraphs[0].font.bold = True
        for i, (_, r) in enumerate(df.iterrows(), start=1):
            for j, c in enumerate(df.columns):
                v = r[c]
                cell = tbl.cell(i, j)
                cell.text = fmt.format(v) if isinstance(v, float) else str(v)
                cell.text_frame.paragraphs[0].font.size = Pt(size)
        return tbl

    def stat(self, s: Any, value: str, label: str, x: int, y: int, w: int = STAT_WIDTH) -> None:
        self.text(s, value, x, y, w, Inches(0.6), size=28, bold=True, color=ACCENT)
        self.text(s, label, x, y + Inches(0.6), w, Inches(0.5), size=11, color=MUTED)

    def flow(self, s: Any, steps: list[tuple[str, str]], x: int, y: int, w: int, box_h: int = FLOW_BOX_HEIGHT) -> None:
        n = len(steps)
        gap = Inches(0.25)
        bw = (w - gap * (n - 1)) / n
        for i, (head, sub) in enumerate(steps):
            bx = x + i * (bw + gap)
            shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y, bw, box_h)
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xED)
            shp.line.color.rgb = LINE
            tf = shp.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = head
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p2 = tf.add_paragraph()
            p2.text = sub
            p2.font.size = Pt(9)
            p2.font.color.rgb = INK
            if i < n - 1:
                arr = s.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    bx + bw + Emu(int(gap * 0.15)),
                    y + box_h / 2 - Inches(0.1),
                    Emu(int(gap * 0.7)),
                    Inches(0.2),
                )
                arr.fill.solid()
                arr.fill.fore_color.rgb = MUTED
                arr.line.fill.background()

    def save(self, path: Path) -> None:
        self.prs.save(str(path))


def _pct(x: float) -> str:
    return f"{100 * x:.0f}%"


def run(cfg: Config) -> None:
    p = cfg.paths
    t, f = p.tables_dir, p.figures_dir
    series = pd.read_parquet(p.interim_dir / "series.parquet")
    overall = pd.read_csv(t / "backtest_metrics_overall.csv").set_index("model")
    hold = pd.read_csv(t / "backtest_metrics_holdout_calibrated_overall.csv").set_index("model")
    risk_methods = pd.read_csv(t / "risk_eval_methods.csv")
    alerts = pd.read_csv(p.output_dir / f"risk_alerts_{cfg.data.as_of}.csv")
    port = pd.read_csv(t / "order_summary_portfolio.csv").iloc[0]
    orders = pd.read_csv(p.output_dir / f"supply_order_{cfg.data.as_of}.csv")
    scored = pd.read_parquet(p.interim_dir / "risk_scored_window.parquet")
    lgbm = overall.loc["lgbm"]
    ma4 = overall.loc["ma4"]
    n_event_series = scored[scored["event"]].groupby(["upc", "cedis"]).ngroups

    d = Deck()

    # 1 title
    s = d.slide("Demand forecasting & supply order generation", "CPG beverages portfolio - 67 SKUs x 6 distribution centers")
    d.text(
        s,
        "A weekly pipeline that forecasts demand with calibrated uncertainty, flags stock-out risk, and recommends "
        "the next replenishment order per SKU and DC - with the business case against today's naive policy.",
        Inches(0.6),
        Inches(1.7),
        Inches(8),
        Inches(1.2),
        size=15,
        color=MUTED,
    )
    d.stat(
        s,
        f"{100 * (1 - lgbm['wape'] / ma4['wape']):.0f}%",
        "lower forecast error than a 4-week moving average",
        Inches(0.6),
        Inches(3.4),
    )
    d.stat(
        s,
        _pct(port["service_level_weighted"]),
        f"implied service level of the recommended order (MA policy: {_pct(port['ma4_service_level_weighted'])})",
        Inches(3.6),
        Inches(3.4),
        Inches(3.2),
    )
    d.stat(
        s,
        f"{int(port['ma4_lines_below_90'])} -> {int(port['lines_below_90'])}",
        "SKU x DC lines below 90% service level (moving-average policy -> recommended)",
        Inches(7.2),
        Inches(3.4),
        Inches(3.0),
    )
    d.stat(
        s,
        f"{port['target_stock_value_freed_by_fq'] / 1e6:,.1f}M",
        f"MXN of target stock ({_pct(port['target_stock_value_freed_by_fq'] / port['target_stock_value'])}) freed by forecast-sized safety stock at the same service targets",
        Inches(10.2),
        Inches(3.4),
        Inches(3.0),
    )
    d.text(
        s,
        "Tracks delivered: A (forecasting) + B (stock-out risk) + supply-order engine. One command re-runs everything: .\\run.ps1",
        Inches(0.6),
        Inches(5.6),
        Inches(12),
        Inches(0.6),
        size=12,
        color=MUTED,
    )

    # 2 data
    s = d.slide("Data: two years of sell-out, three weeks of inventory", "Data analysis & preparation")
    d.bullets(
        s,
        [
            f"{len(series)} SKU x DC series, daily sell-out 2024-03-18 to 2026-04-10, modelled on 107 complete weeks.",
            "Store inventory exists only for 2026-03-20 .. 04-09; the last complete snapshot (2026-04-02) is 'current on-hand'. The final week has 4 UPCs.",
            "Clean-ups: 3 UPCs mapped to two item numbers (de-duplicated), 7,116 negative store on-hand rows (clipped, kept as a stock-out signal), late-starting series kept as missing, 1,720 daily outliers winsorised for training only.",
            "Calendar effects engineered: Mexican federal holidays, quincena paydays, Semana Santa, El Buen Fin, December peak. Future price and promos are treated as unknown.",
            f"Clusters (ABC x XYZ): {', '.join(f'{k} {v}' for k, v in series['cluster'].value_counts().sort_index().items())}. Lead time is 7 days and MOQ 100 for every SKU; only safety-stock days vary (3 / 7 / 14).",
        ],
        Inches(0.6),
        Inches(1.6),
        Inches(7.2),
        Inches(5),
        size=12,
    )
    d.picture(s, f / "inventory_coverage.png", Inches(8.0), Inches(1.7), w=Inches(4.9))
    d.picture(s, f / "portfolio_forecast.png", Inches(8.0), Inches(4.2), w=Inches(4.9))

    # 3 architecture
    s = d.slide("Architecture and data flow", "How it runs")
    d.flow(
        s,
        [
            ("Extraction", "sell-out, inventory, catalogs; schema validation"),
            ("Feature store", "daily + weekly parquet: lags, calendar, promo, cover"),
            ("Training", "expanding-window backtest; per-cluster selection; interval calibration"),
            ("Scoring", "8-week quantile forecast per SKU x DC, weekly"),
            ("Recommendation", "risk scores + order-up-to policy -> order file"),
            ("Consumption", "planner queue, S&OP input, API"),
        ],
        Inches(0.6),
        Inches(1.8),
        Inches(12.1),
        Inches(1.4),
    )
    d.bullets(
        s,
        [
            "Stages communicate only through files (parquet / CSV), so any stage can be re-run alone and every artifact is inspectable.",
            "All assumptions live in config.toml: as-of date, cold-start threshold, outlier z, stock-out label share, quantiles, folds, service-level targets, cost ratio.",
            "Typed Python package, argparse CLI, pytest suite (metrics, calendar, aggregation, MOQ rounding, service-level maths, labels).",
            "Runtime on a laptop: about 3 minutes end to end (LightGBM global model + 328 ETS fits x 8 folds).",
        ],
        Inches(0.6),
        Inches(3.6),
        Inches(12),
        Inches(3),
        size=12,
    )

    # 4 forecasting results
    s = d.slide("Forecasting: a global LightGBM model wins every cluster", "Track A - results")
    d.picture(s, f / "model_comparison.png", Inches(0.5), Inches(1.6), w=Inches(7.3))
    tbl = overall.reset_index()[["model", "wape", "bias", "coverage_90"]].rename(columns={"coverage_90": "cov90 raw"})
    d.table(s, tbl, Inches(8.1), Inches(1.7), Inches(4.8), Inches(1.6), size=10)
    d.bullets(
        s,
        [
            f"Overall WAPE {lgbm['wape']:.2f} vs {overall.loc['ets', 'wape']:.2f} (ETS) and {ma4['wape']:.2f} (MA4); seasonal naive {overall.loc['seasonal_naive', 'wape']:.2f} with +{100 * overall.loc['seasonal_naive', 'bias']:.0f}% bias - last year is not a guide.",
            "Selection rule: lowest mean WAPE over 8 folds unless a simpler or steadier model is within 0.02. LightGBM wins outright in all six clusters.",
            f"Bias {lgbm['bias']:+.2f}: slight under-forecast, concentrated in promo weeks.",
        ],
        Inches(8.1),
        Inches(3.5),
        Inches(4.8),
        Inches(3.3),
        size=11,
    )

    # 5 uncertainty
    s = d.slide("Uncertainty: calibrated intervals, honestly measured", "Track A - prediction intervals")
    d.picture(s, f / "coverage_calibration.png", Inches(0.5), Inches(1.6), w=Inches(6.2))
    d.picture(s, f / "cluster_forecast_vs_actual.png", Inches(6.9), Inches(1.6), w=Inches(6.2))
    d.bullets(
        s,
        [
            f"Raw LightGBM quantiles were over-confident (90% band covered {_pct(lgbm['coverage_90'])}). A conformal-style width factor per model x cluster is fitted on folds 1-5 and checked on folds 6-8: calibrated 90% coverage {_pct(hold.loc['lgbm', 'coverage_90'])}, 80% coverage {_pct(hold.loc['lgbm', 'coverage_80'])}.",
            "The calibrated p90 is what sizes safety stock in the order engine; the median drives the expected quantity.",
            "Error grows with horizon (see report); the 2-week protection period uses the most accurate part of the curve.",
        ],
        Inches(0.6),
        Inches(4.9),
        Inches(12),
        Inches(2.2),
        size=11,
    )

    # 6 risk
    s = d.slide("Stock-out risk: an alert planners can read", "Track B - results")
    d.bullets(
        s,
        [
            "Alert = 'on current DC stock and expected demand, this SKU is likely short at stores within the 7-day lead time'. High severity: P(stock-out) > 75% or cover below lead time.",
            "Three scorers: days-of-cover rule (cover < lead time + safety days), forecast probability P(7-day demand > on-hand) from calibrated quantiles, Isolation Forest on cover / probability / sales-vs-forecast / stock trend.",
            f"Label: >= {_pct(cfg.data.stockout_store_share)} of a DC's stores at zero. The window holds {int(scored['event'].sum())} event days over {n_event_series} series - all chronic, so lead-time-to-alert cannot be measured from this file.",
            f"As of {cfg.data.as_of}: {alerts['severity'].value_counts().to_dict()} alerts across {len(alerts)} SKU x DC.",
        ],
        Inches(0.6),
        Inches(1.6),
        Inches(6.2),
        Inches(4.5),
        size=11,
    )
    d.table(
        s,
        risk_methods[["method", "n_alerts", "n_events", "precision", "recall", "false_alarm_rate"]],
        Inches(0.6),
        Inches(5.2),
        Inches(6.2),
        Inches(1.4),
        size=10,
    )
    d.picture(s, f / "risk_window.png", Inches(7.0), Inches(1.6), w=Inches(6.0))

    # 7 orders
    s = d.slide("Supply order: what to buy next cycle, and what it costs", "Recommendation engine")
    d.picture(s, f / "orders_summary.png", Inches(0.5), Inches(1.6), w=Inches(7.2))
    d.bullets(
        s,
        [
            "Order-up-to over a 14-day protection period (lead time 7 + weekly review 7). Safety stock = max(catalog days policy, stock needed for the ABC service-level target under the calibrated forecast). Rounded up to MOQ.",
            f"Recommended: {port['order_units']:,.0f} units, MXN {port['working_capital']:,.0f} working capital, implied service level {_pct(port['service_level_weighted'])}, expected fill rate {_pct(port['fill_rate'])}.",
            f"4-week moving-average policy: {port['ma4_order_units']:,.0f} units, service level {_pct(port['ma4_service_level_weighted'])}, fill rate {_pct(port['ma4_fill_rate'])}.",
            f"Delta this cycle: {port['delta_working_capital_vs_ma4']:+,.0f} MXN working capital for {port['lost_units_reduction_vs_ma4']:,.0f} fewer expected lost units; the tail matters more: {int(port['ma4_lines_below_90'])} -> {int(port['lines_below_90'])} lines below 90% service level.",
            f"Forecast-sized safety stock at the same targets (fq variant) would carry {port['fq_target_stock_value'] / 1e6:,.1f}M instead of {port['target_stock_value'] / 1e6:,.1f}M MXN of target stock at {_pct(port['fq_service_level_weighted'])} service level - the negotiable lever.",
            f"{int((orders['flags'].fillna('') != '').sum())} of {len(orders)} lines flagged for planner review (discontinued, no inventory snapshot, high promo).",
        ],
        Inches(8.0),
        Inches(1.6),
        Inches(5.0),
        Inches(5.2),
        size=10,
    )

    # 8 impact + trust
    s = d.slide("Business impact and where not to trust it", "Framing")
    d.bullets(
        s,
        [
            f"Fill-rate lift: {_pct(port['ma4_fill_rate'])} -> {_pct(port['fill_rate'])} over the protection period (expected, model distribution).",
            f"Lost-sales reduction: {port['lost_units_reduction_vs_ma4']:,.0f} units per cycle at portfolio level.",
            f"Inventory cash: {port['target_stock_value_freed_by_fq'] / 1e6:,.1f}M MXN of target stock freed if safety stock is sized from the calibrated forecast instead of fixed days (shelf-price basis; set cost_ratio for cost).",
            f"Service-level tail: {int(port['ma4_lines_below_90'])} lines below 90% under the moving-average policy vs {int(port['lines_below_90'])} recommended.",
        ],
        Inches(0.6),
        Inches(1.6),
        Inches(6),
        Inches(3),
        size=12,
    )
    d.text(s, "Do not trust the model for", Inches(7.0), Inches(1.6), Inches(6), Inches(0.4), size=14, bold=True, color=GOLD)
    d.bullets(
        s,
        [
            f"Discontinued lines ({int(series['is_discontinued'].sum())} series): orders forced to zero, flagged.",
            "Cold-start SKUs (< 26 weeks): moving-average fallback, flagged.",
            "Promo-driven weeks: promos are unknown ahead, spikes are under-forecast - overlay the promo calendar.",
            "A-Z cluster: highest error and widest bands - order on the p90, not the median.",
            "Lines without an inventory snapshot or with in-transit orders (receipts are not in the data).",
        ],
        Inches(7.0),
        Inches(2.1),
        Inches(6),
        Inches(4.5),
        size=11,
    )
    d.picture(s, f / "series_examples.png", Inches(0.6), Inches(4.2), w=Inches(6.0))

    # 9 productionization
    s = d.slide("Productionization plan", "Operating the pipeline")
    d.flow(
        s,
        [
            ("Schedule", "weekly scoring Monday after sales close; monthly LightGBM retrain; ETS refit each run"),
            (
                "Monitoring",
                "rolling 4-wk WAPE & bias per cluster vs backtest envelope; interval coverage vs nominal; input schema & coverage checks",
            ),
            ("Alerting", "drift breach -> retrain + human review; data gaps -> hold orders for affected DCs"),
            ("Human in the loop", "planner approves / overrides each order line; overrides logged as training signal"),
            ("Ownership", "DS owns models & monitors; supply planning owns policy targets in config.toml"),
        ],
        Inches(0.6),
        Inches(1.8),
        Inches(12.1),
        Inches(1.9),
    )
    d.bullets(
        s,
        [
            "Failure modes covered: missing inventory day (falls back to last full snapshot and flags), new SKU (cold-start rule), catalog change (schema validation fails loudly), negative stock (clipped and counted).",
            "Retraining trigger: cluster WAPE above the backtest p90 for two consecutive weeks, or coverage below nominal minus 10 points.",
            "Reproducibility: pinned dependencies, fixed seeds, one command from a clean clone, tests on every metric and policy rule.",
        ],
        Inches(0.6),
        Inches(4.1),
        Inches(12),
        Inches(3),
        size=12,
    )

    # 10 next steps
    s = d.slide("Next steps and integration", "Roadmap")
    d.bullets(
        s,
        [
            "Planner UI: order lines with severity, cover days, p50/p90 demand and the override reason; one click to approve a DC.",
            "S&OP cycle: monthly cluster-level forecast pack (this report) as the statistical baseline; commercial overlays for promos and launches.",
            "Data asks: promo calendar and planned prices (largest accuracy lever), receipts / in-transit orders, store-level stock history beyond 21 days to learn true stock-out dynamics.",
            "Model roadmap: promo-aware features once the calendar exists; hierarchical reconciliation UPC x DC -> UPC; store-level allocation of the DC order.",
            "API: score-on-demand endpoint returning forecast, risk and recommended order for a SKU x DC; batch file for ERP upload.",
        ],
        Inches(0.6),
        Inches(1.6),
        Inches(12),
        Inches(5),
        size=13,
    )

    out = p.reports_dir / "deck.pptx"
    d.save(out)
    log.info("wrote %s (%d slides)", out, d.n)
