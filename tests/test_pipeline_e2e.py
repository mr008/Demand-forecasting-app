"""Comprehensive test file: unit tests for the evaluation gates, stage-level input/output
tests, and a full end-to-end run over a synthetic dataset with known ground truth.

Layout
------
1. Unit: ``Gate`` comparison semantics and the synthetic generator itself.
2. Stage I/O (on the session-wide synthetic run): what each stage must write, and
   the invariants those files must satisfy (schemas, keys, monotone quantiles,
   MOQ multiples, flags, labels, onsets).
3. Ground truth: forecasts vs the planted demand process; discontinued and
   cold-start handling; the stock-out we planted is alerted; determinism on re-run.
4. The quality-gate scorecard over the synthetic run must have no hard failures.

Run only the fast tests with ``pytest -m "not e2e"``.
"""

from __future__ import annotations

import hashlib

import numpy as np
import pandas as pd
import pytest
from pptx import Presentation

from supply_pipeline import cli, evaluate
from supply_pipeline.evaluate import Gate, run_gates
from supply_pipeline.models import MODEL_ORDER
from tests.synthetic import (
    CEDIS_MULT,
    COLD_START_SPEC,
    DISCONTINUED_SPEC,
    DUAL_PRIME_SPEC,
    LAST_FULL_INV_DAY,
    ONSET_DAY,
    ONSET_SPEC,
    SPECS,
    STEADY_SPEC,
    STOCKOUT_SPEC,
    SyntheticDataset,
)

SERIES_KEY = ["upc", "cedis"]


# =============================================================================
# 1. Unit tests
# =============================================================================
class TestGate:
    def test_ops(self) -> None:
        assert Gate("a", "x", 0.1, "<=", 0.2, "hard").passed
        assert not Gate("a", "x", 0.3, "<=", 0.2, "hard").passed
        assert Gate("a", "x", 0.5, ">=", 0.5, "soft").passed
        assert Gate("a", "x", 3.0, "==", 3.0, "hard").passed
        assert Gate("a", "x", 0.9, "in", (0.85, 0.97), "hard").passed
        assert not Gate("a", "x", 0.98, "in", (0.85, 0.97), "hard").passed

    def test_nan_never_passes(self) -> None:
        assert not Gate("a", "x", float("nan"), ">=", 0.0, "soft").passed

    def test_status_reflects_severity(self) -> None:
        assert Gate("a", "x", 1.0, "<=", 0.5, "hard").status == "FAIL"
        assert Gate("a", "x", 1.0, "<=", 0.5, "soft").status == "warn"
        assert Gate("a", "x", 0.1, "<=", 0.5, "soft").status == "pass"

    def test_row_formats_range_threshold(self) -> None:
        row = Gate("a", "x", 0.9, "in", (0.85, 0.97), "hard", "note").row()
        assert row["threshold"] == "[0.85, 0.97]" and row["note"] == "note"


class TestSyntheticDataset:
    def test_raw_files_and_planted_messiness(self, synthetic: SyntheticDataset) -> None:
        raw = synthetic.root / "data" / "raw"
        sell = pd.read_csv(raw / "challenge_daily_sell_out_pricing.csv", parse_dates=["date"])
        inv = pd.read_csv(raw / "challenge_inventory.csv", parse_dates=["date"])
        stores = pd.read_csv(raw / "challenge_store_catalog.csv")
        cat = pd.read_csv(raw / "challenge_upc_catalog.csv")

        assert sell.groupby(SERIES_KEY).ngroups == synthetic.n_series == 23
        assert not sell.duplicated(["date", "upc", "cedis"]).any()
        assert sell["date"].min() == pd.Timestamp("2024-03-18") and sell["date"].max() == pd.Timestamp("2026-04-10")
        assert 0.05 < sell["promo_flag"].mean() < 0.25
        assert (sell["sell_out_pzs"] % 1 != 0).mean() > 0.5  # fractional like the real file
        late = sell[sell["upc"] == COLD_START_SPEC.upc]["date"].min()
        assert late == COLD_START_SPEC.start

        assert len(cat) == len(SPECS) + 1 and cat["upc"].nunique() == len(SPECS)  # dual prime item
        assert (inv["on_hand_qty"] < 0).sum() > 0
        assert inv.duplicated(["date", "upc", "store_nbr"]).sum() > 0  # dual prime rows
        assert inv[inv["date"] > LAST_FULL_INV_DAY]["upc"].nunique() == 2
        assert set(stores["cedis"]) == set(CEDIS_MULT) and stores["store_name"].isna().sum() == 1

    def test_truth_is_deterministic_process(self, synthetic: SyntheticDataset) -> None:
        t = synthetic.truth
        steady = t[(t["upc"] == STEADY_SPEC.upc) & (t["cedis"] == "X1")]
        # season (0.8-1.2) x weekday (0.9-1.25) x payday (1-1.15) bounds the deterministic process
        assert steady["true_mean"].between(0.55 * STEADY_SPEC.level, 1.8 * STEADY_SPEC.level).all()
        disc = t[(t["upc"] == DISCONTINUED_SPEC.upc) & (t["date"] >= "2026-03-01")]
        assert (disc["true_mean"] < 0.1 * DISCONTINUED_SPEC.level).all()
        assert t["date"].max() > pd.Timestamp("2026-04-10")  # truth covers the forecast horizon


# =============================================================================
# 2. Stage input / output tests (session-wide synthetic run)
# =============================================================================
@pytest.mark.e2e
class TestPrepareStage:
    def test_artifacts_exist(self, pipeline_run) -> None:
        p = pipeline_run.paths
        for name in ("catalog", "stores", "daily", "inventory_cedis_daily", "weekly", "series", "calendar"):
            assert (p.interim_dir / f"{name}.parquet").exists(), name
        assert (p.tables_dir / "coverage_inventory_by_date.csv").exists()

    def test_catalog_deduped_and_clusters(self, pipeline_run) -> None:
        cat = pd.read_parquet(pipeline_run.paths.interim_dir / "catalog.parquet")
        assert len(cat) == len(SPECS) and cat["upc"].is_unique
        assert set(cat["cluster"]) == {s.cluster for s in SPECS}

    def test_weekly_table_shape_and_reconciliation(self, pipeline_run, synthetic: SyntheticDataset) -> None:
        w = pd.read_parquet(pipeline_run.paths.interim_dir / "weekly.parquet")
        d = pd.read_parquet(pipeline_run.paths.interim_dir / "daily.parquet")
        assert w.groupby(SERIES_KEY).ngroups == synthetic.n_series
        assert w["week_start"].nunique() == 107
        assert (w["week_start"].dt.dayofweek == 0).all()
        assert not w.duplicated(SERIES_KEY + ["week_start"]).any()
        # Full weeks reconcile exactly to the daily file.
        d["week_start"] = (d["date"] - pd.to_timedelta(d["date"].dt.dayofweek, unit="D")).dt.normalize()
        ds = d.groupby(SERIES_KEY + ["week_start"])["sell_out_pzs"].agg(["sum", "count"]).reset_index()
        m = w.merge(ds, on=SERIES_KEY + ["week_start"])
        full = m[m["count"] == 7]
        assert np.allclose(full["y"], full["sum"])
        # Calendar counts are populated and bounded.
        assert w["payday_days"].between(0, 3).all() and w["holiday_days"].between(0, 7).all()

    def test_series_flags(self, pipeline_run) -> None:
        s = pd.read_parquet(pipeline_run.paths.interim_dir / "series.parquet").set_index(SERIES_KEY)
        assert s.loc[(DISCONTINUED_SPEC.upc, "X1"), "is_discontinued"]
        assert s.loc[(COLD_START_SPEC.upc, "X1"), "is_cold_start"]
        assert not s.loc[(STEADY_SPEC.upc, "X1"), "is_cold_start"]
        assert s.loc[(STEADY_SPEC.upc, "X1"), "n_weeks"] == 107
        assert s["has_inventory"].all()
        assert (s["interior_gap_days"] == 0).all()

    def test_inventory_aggregation_keeps_stockout_signal(self, pipeline_run) -> None:
        inv = pd.read_parquet(pipeline_run.paths.interim_dir / "inventory_cedis_daily.parquet")
        assert (inv["on_hand"] >= 0).all()  # clipped
        assert (inv["stores_negative"] > 0).any()  # ...but negatives were counted
        chronic = inv[(inv["upc"] == STOCKOUT_SPEC.upc) & (inv["cedis"] == "X1")]
        assert (chronic["stockout_store_share"] >= 0.5).all()
        onset = inv[(inv["upc"] == ONSET_SPEC.upc) & (inv["cedis"] == "X2")].set_index("date")["stockout_store_share"]
        assert (onset[onset.index < ONSET_DAY] < 0.25).all()
        assert (onset[onset.index >= ONSET_DAY] >= 0.4).all()
        # Dual prime-item rows were summed, not duplicated, per store.
        dual = inv[inv["upc"] == DUAL_PRIME_SPEC.upc]
        assert (dual["stores_reporting"] <= 10).all()


@pytest.mark.e2e
class TestBacktestStage:
    def test_predictions_cover_models_and_folds(self, pipeline_run) -> None:
        pred = pd.read_parquet(pipeline_run.paths.interim_dir / "backtest_predictions.parquet")
        assert set(pred["model"]) == set(MODEL_ORDER)
        assert pred["fold"].nunique() == pipeline_run.forecast.backtest_folds
        assert pred["h"].max() == pipeline_run.forecast.horizon_weeks
        assert pred["y"].notna().mean() > 0.95  # actuals joined for the evaluated horizon
        # Cold-start series is excluded from every backtest origin (too little history).
        assert COLD_START_SPEC.upc not in set(pred["upc"])

    def test_metric_tables_and_selection(self, pipeline_run) -> None:
        t = pipeline_run.paths.tables_dir
        cluster = pd.read_csv(t / "backtest_metrics_cluster.csv")
        sel = pd.read_csv(t / "model_selection.csv")
        calib = pd.read_csv(t / "interval_calibration.csv")
        assert set(cluster["cluster"]) == {s.cluster for s in SPECS}
        assert sel["cluster"].is_unique and set(sel["cluster"]) == set(cluster["cluster"])
        assert sel["selected_model"].isin(MODEL_ORDER).all()
        # WAPE is unbounded above (seasonal naive on a discontinued cluster is legitimately huge), so only check it is finite.
        assert np.isfinite(cluster["wape"]).all() and (cluster["wape"] >= 0).all()
        assert cluster["coverage_90"].between(0, 1).all()
        assert len(calib) == 2 * cluster["model"].nunique() * cluster["cluster"].nunique()
        assert (calib["k"] >= -0.4).all()
        for name in (
            "backtest_metrics_fold",
            "backtest_metrics_cluster_h",
            "backtest_metrics_overall",
            "backtest_metrics_cedis",
            "backtest_metrics_holdout_raw",
            "backtest_metrics_holdout_calibrated",
        ):
            assert (t / f"{name}.csv").exists(), name

    def test_calibration_improves_or_keeps_coverage_of_overconfident_models(self, pipeline_run) -> None:
        t = pipeline_run.paths.tables_dir
        raw = pd.read_csv(t / "backtest_metrics_holdout_raw.csv").groupby("model")["coverage_90"].mean()
        cal = pd.read_csv(t / "backtest_metrics_holdout_calibrated.csv").groupby("model")["coverage_90"].mean()
        for m in raw.index:
            if raw[m] < 0.85:
                assert cal[m] > raw[m], m


@pytest.fixture(scope="module")
def final(pipeline_run) -> pd.DataFrame:
    return pd.read_csv(
        pipeline_run.paths.output_dir / f"forecast_{pipeline_run.data.last_complete_week_start}.csv",
        parse_dates=["target_week", "origin"],
    )


@pytest.fixture(scope="module")
def orders(pipeline_run) -> pd.DataFrame:
    o = pd.read_csv(pipeline_run.paths.output_dir / f"supply_order_{pipeline_run.data.as_of}.csv")
    o["flags"] = o["flags"].fillna("")
    return o


@pytest.mark.e2e
class TestForecastStage:
    def test_both_origins_written(self, pipeline_run) -> None:
        p = pipeline_run.paths
        assert (p.output_dir / f"forecast_{pipeline_run.data.last_complete_week_start}.csv").exists()
        assert (p.output_dir / f"forecast_{pipeline_run.risk.eval_origin}.csv").exists()
        assert (p.interim_dir / "forecasts_all_models.parquet").exists()

    def test_schema_and_invariants(self, final: pd.DataFrame, pipeline_run, synthetic: SyntheticDataset) -> None:
        q = ["q05", "q10", "q50", "q90", "q95"]
        assert {*SERIES_KEY, "h", "target_week", "model", "cluster", "is_cold_start", "is_discontinued", *q} <= set(final.columns)
        arr = final[q].to_numpy()
        assert np.isfinite(arr).all() and (arr >= 0).all()
        assert (np.diff(arr, axis=1) >= -1e-9).all()
        assert (final.groupby(SERIES_KEY)["h"].nunique() == pipeline_run.forecast.horizon_weeks).all()
        assert final.groupby(SERIES_KEY).ngroups == synthetic.n_series
        assert (final.groupby(SERIES_KEY)["model"].nunique() == 1).all()
        origin = pd.Timestamp(pipeline_run.data.last_complete_week_start)
        assert (final["target_week"] == origin + pd.to_timedelta(final["h"] * 7, unit="D")).all()

    def test_cold_start_uses_fallback_model(self, final: pd.DataFrame) -> None:
        cs = final[final["upc"] == COLD_START_SPEC.upc]
        assert cs["is_cold_start"].all() and (cs["model"] == "ma4").all()
        others = final[final["upc"] != COLD_START_SPEC.upc]
        assert not others["is_cold_start"].any()

    def test_discontinued_forecasts_near_zero(self, final: pd.DataFrame, pipeline_run) -> None:
        s = pd.read_parquet(pipeline_run.paths.interim_dir / "series.parquet").set_index(SERIES_KEY)
        d = final[(final["upc"] == DISCONTINUED_SPEC.upc) & (final["cedis"] == "X1")]
        assert d["is_discontinued"].all()
        assert d["q50"].mean() < 0.1 * s.loc[(DISCONTINUED_SPEC.upc, "X1"), "mean_weekly"]


@pytest.mark.e2e
class TestRiskStage:
    def test_tables_and_window_labels(self, pipeline_run) -> None:
        p = pipeline_run.paths
        methods = pd.read_csv(p.tables_dir / "risk_eval_methods.csv")
        sweep = pd.read_csv(p.tables_dir / "risk_threshold_sweep.csv")
        lead = pd.read_csv(p.tables_dir / "risk_lead_time.csv")
        scored = pd.read_parquet(p.interim_dir / "risk_scored_window.parquet")
        assert set(methods["method"]) == {"cover", "prob", "iforest"}
        assert methods[["precision", "recall", "false_alarm_rate"]].fillna(0).map(lambda v: 0 <= v <= 1).all().all()
        assert set(sweep["method"]) == {"cover", "prob"}
        assert scored["p_stockout_7d"].between(0, 1).all()
        # Planted chronic stock-out is an event on every day; the planted onset is detected as an onset.
        chronic = scored[(scored["upc"] == STOCKOUT_SPEC.upc) & (scored["cedis"] == "X1")]
        assert chronic["event"].all()
        onset = scored[(scored["upc"] == ONSET_SPEC.upc) & (scored["cedis"] == "X2") & scored["episode_onset"]]
        assert list(onset["date"]) == [ONSET_DAY]
        assert int(lead["n_onsets"].iloc[0]) >= 1
        # Labels are NaN exactly where the 7-day lookahead is not fully observed.
        last_obs = scored.groupby(SERIES_KEY)["date"].transform("max")
        assert scored.loc[scored["date"] > last_obs - pd.Timedelta(days=7), "label_7d"].isna().all()

    def test_as_of_alerts(self, pipeline_run, synthetic: SyntheticDataset) -> None:
        alerts = pd.read_csv(pipeline_run.paths.output_dir / f"risk_alerts_{pipeline_run.data.as_of}.csv")
        assert len(alerts) == synthetic.n_series
        assert set(alerts["severity"]) <= {"none", "medium", "high"}
        planted = alerts[(alerts["upc"] == STOCKOUT_SPEC.upc) & (alerts["cedis"] == "X1")].iloc[0]
        assert planted["severity"] != "none"
        assert planted["alert_cover"] or planted["alert_prob"]
        assert planted["cover_days"] < planted["cover_threshold"]


@pytest.mark.e2e
class TestOrdersStage:
    def test_one_line_per_series_and_policy_invariants(self, orders: pd.DataFrame, synthetic: SyntheticDataset) -> None:
        assert len(orders) == synthetic.n_series and not orders.duplicated(SERIES_KEY).any()
        assert (orders["order_qty"] >= 0).all()
        assert (((orders["order_qty"] % orders["moq"]).abs() < 1e-6) | (orders["order_qty"] == 0)).all()
        assert (orders.loc[orders["upc"] == SPECS[7].upc, "moq"] == 50).all()
        # Recompute the rounding from the reported components.
        raw = (orders["order_up_to"] - orders["on_hand_projected"]).clip(lower=0)
        expected = np.where(raw <= 0, 0, np.ceil(raw / orders["moq"]) * orders["moq"])
        expected = np.where(orders["flags"].str.contains("discontinued"), 0, expected)
        assert np.allclose(orders["order_qty"], expected)
        assert np.allclose(orders["working_capital"], orders["order_qty"] * orders["unit_value"])
        assert orders["implied_service_level"].between(0, 1).all()
        assert (orders["safety_stock"] >= orders["safety_stock_policy"] - 1e-6).all()
        assert orders["safety_stock_binding"].isin(["policy_days", "forecast_quantile"]).all()

    def test_flags(self, orders: pd.DataFrame) -> None:
        disc = orders[orders["upc"] == DISCONTINUED_SPEC.upc]
        assert disc["flags"].str.contains("discontinued_review").all() and (disc["order_qty"] == 0).all()
        cold = orders[orders["upc"] == COLD_START_SPEC.upc]
        assert cold["flags"].str.contains("cold_start").all()
        assert (cold["model"] == "ma4").all()

    def test_stocked_out_dc_gets_an_order(self, orders: pd.DataFrame) -> None:
        line = orders[(orders["upc"] == STOCKOUT_SPEC.upc) & (orders["cedis"] == "X1")].iloc[0]
        assert line["order_qty"] > 0
        assert line["implied_service_level"] >= line["service_level_target"] - 1e-6
        # Half the stores are empty, so the DC holds about half its normal stock: the order must at least
        # restore the order-up-to level.
        assert line["on_hand_projected"] + line["order_qty"] >= line["order_up_to"] - 1e-6

    def test_unit_value_tracks_shelf_price(self, orders: pd.DataFrame) -> None:
        for spec in SPECS:
            lines = orders[orders["upc"] == spec.upc]
            assert lines["unit_value"].between(0.7 * spec.base_price, 1.05 * spec.base_price).all(), spec.upc

    def test_summary_tables(self, pipeline_run, orders: pd.DataFrame) -> None:
        t = pipeline_run.paths.tables_dir
        port = pd.read_csv(t / "order_summary_portfolio.csv").iloc[0]
        by_cluster = pd.read_csv(t / "order_summary_cluster.csv")
        assert port["n_series"] == len(orders)
        assert np.isclose(port["order_units"], orders["order_qty"].sum())
        assert np.isclose(by_cluster["working_capital"].sum(), orders["working_capital"].sum())
        assert 0 <= port["service_level_weighted"] <= 1 and 0 <= port["fq_service_level_weighted"] <= 1
        assert port["fq_target_stock_value"] <= port["target_stock_value"] + 1e-6


@pytest.mark.e2e
class TestReportAndDeck:
    def test_report_outputs(self, pipeline_run) -> None:
        p = pipeline_run.paths
        text = (p.reports_dir / "summary.md").read_text(encoding="utf-8")
        for section in evaluate.REQUIRED_SUMMARY_SECTIONS:
            assert section in text
        assert "nan" not in text.lower().replace("nanotech", "")  # no unformatted NaN leaked into the narrative
        assert len(list(p.figures_dir.glob("*.png"))) >= evaluate.EXPECTED_FIGURES

    def test_deck(self, pipeline_run) -> None:
        deck = Presentation(str(pipeline_run.paths.reports_dir / "deck.pptx"))
        assert len(deck.slides) == evaluate.EXPECTED_SLIDES
        titles = [sh.text_frame.text for s in deck.slides for sh in s.shapes if sh.has_text_frame and sh.text_frame.text]
        assert any("Supply order" in t or "supply order" in t for t in titles)


# =============================================================================
# 3. Ground-truth and behavioural tests
# =============================================================================
@pytest.mark.e2e
class TestGroundTruth:
    def test_steady_series_forecast_tracks_true_process(self, pipeline_run, synthetic: SyntheticDataset) -> None:
        fc = pd.read_csv(
            pipeline_run.paths.output_dir / f"forecast_{pipeline_run.data.last_complete_week_start}.csv",
            parse_dates=["target_week"],
        )
        t = synthetic.truth.copy()
        t["week_start"] = (t["date"] - pd.to_timedelta(t["date"].dt.dayofweek, unit="D")).dt.normalize()
        truth_w = t.groupby(SERIES_KEY + ["week_start"])["true_mean"].sum().rename("truth").reset_index()
        f = fc[(fc["upc"] == STEADY_SPEC.upc) & (fc["cedis"] == "X1")].merge(
            truth_w.rename(columns={"week_start": "target_week"}), on=SERIES_KEY + ["target_week"]
        )
        assert len(f) == pipeline_run.forecast.horizon_weeks
        mape = (f["q50"] - f["truth"]).abs().div(f["truth"]).mean()
        assert mape < 0.25, f"steady A-X series median off by {mape:.1%} from the true process"
        inside = ((f["truth"] >= f["q05"]) & (f["truth"] <= f["q95"])).mean()
        assert inside >= 0.75, "calibrated 90% band should contain the true mean in most weeks"

    def test_backtest_beats_naive_on_low_noise_cluster(self, pipeline_run) -> None:
        cluster = pd.read_csv(pipeline_run.paths.tables_dir / "backtest_metrics_cluster.csv")
        ax = cluster[cluster["cluster"] == STEADY_SPEC.cluster].set_index("model")["wape"]
        assert ax.min() <= ax["ma4"] + 1e-9
        assert ax[["ets", "lgbm"]].min() < 0.25

    def test_rerun_of_downstream_stages_is_deterministic(self, pipeline_run, synthetic: SyntheticDataset) -> None:
        p = pipeline_run.paths
        targets = [
            p.output_dir / f"risk_alerts_{pipeline_run.data.as_of}.csv",
            p.output_dir / f"supply_order_{pipeline_run.data.as_of}.csv",
        ]
        before = [hashlib.sha256(x.read_bytes()).hexdigest() for x in targets]
        for stage in ("risk", "orders"):
            assert cli.main(["--config", str(synthetic.config_path), stage]) == 0
        after = [hashlib.sha256(x.read_bytes()).hexdigest() for x in targets]
        assert before == after


# =============================================================================
# 4. Quality gates over the synthetic run
# =============================================================================
@pytest.mark.e2e
def test_eval_scorecard_has_no_hard_failures(pipeline_run) -> None:
    scorecard = run_gates(pipeline_run)
    evaluate.write_scorecard(pipeline_run, scorecard)
    assert (pipeline_run.paths.tables_dir / "eval_scorecard.csv").exists()
    failed = scorecard[scorecard["status"] == "FAIL"]
    assert failed.empty, "hard gates failed:\n" + failed[["area", "gate", "value", "op", "threshold"]].to_string(index=False)
    assert len(scorecard) >= 35


@pytest.mark.e2e
def test_evaluate_stage_exits_non_zero_when_a_hard_gate_fails(pipeline_run, monkeypatch: pytest.MonkeyPatch) -> None:
    def broken(cfg):
        return [Gate("test", "forced", 1.0, "<=", 0.0, "hard")]

    monkeypatch.setattr(evaluate, "report_gates", broken)
    with pytest.raises(evaluate.EvalError):
        evaluate.run(pipeline_run)
    assert (pipeline_run.paths.reports_dir / "eval_report.md").exists()
